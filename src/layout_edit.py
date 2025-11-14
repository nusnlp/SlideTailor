from pptx.util import Pt
from copy import deepcopy
from functools import partial

import matplotlib.pyplot as plt
import matplotlib.patches as patches

from presentation import SlidePage, Picture, TextBox, Closure
from llms import Role


def layout_parsing(page, normalize=False):

    slide_width = page.slide_width
    slide_height = page.slide_height
    
    parsed_shapes = {}
    for shape in page.shapes:
        
        left = shape.style["shape_bounds"]["left"].pt
        top = shape.style["shape_bounds"]["top"].pt
        width = shape.style["shape_bounds"]["width"].pt
        height = shape.style["shape_bounds"]["height"].pt
            
        if normalize:
            left = shape.style["shape_bounds"]["left"].pt / slide_width
            top = shape.style["shape_bounds"]["top"].pt / slide_height
            width = shape.style["shape_bounds"]["width"].pt / slide_width
            height = shape.style["shape_bounds"]["height"].pt / slide_height
        
        parsed_shapes[f"{shape.shape_idx}"] = {
            "type": shape.__class__.__name__,
            "left": left,
            "top": top,
            "width": width,
            "height": height
        }
        
    return parsed_shapes, (slide_width, slide_height)


def draw_bboxes(elements, figsize=(6, 6), edgecolor="tab:red", textcolor="tab:blue", notation="",
                slide_width=720, slide_height=540):
    """
    Draw every bounding box stored in `elements` on a Matplotlib figure.
    """
    fig, ax = plt.subplots(figsize=figsize)


    # Draw each rectangle and its index
    for idx, elem in elements.items():
        x, y, w, h = elem["left"], elem["top"], elem["width"], elem["height"]
        rect = patches.Rectangle((x, y), w, h,
                                linewidth=1.5, edgecolor=edgecolor,
                                facecolor='none')
        ax.add_patch(rect)

        # Put the index slightly above the top-left corner
        ax.text(x, y - 5, str(idx), color=textcolor,
                ha="left", va="top", fontsize=9, weight="bold")
        
        # Add the element type in the center of the bbox
        if "type" in elem:
            # Calculate center of the box
            center_x = x + w/2
            center_y = y + h/2
            # Add type text
            ax.text(center_x, center_y, elem["type"], color=textcolor,
                    ha="center", va="center", fontsize=8)

    # Make the origin the top-left like typical UI coords
    ax.set_xlim(0, slide_width)
    ax.set_ylim(slide_height, 0)          # invert y-axis
    ax.set_aspect("equal", adjustable="box")
    ax.axis("off")                 # hide axes
    
    if notation:
        ax.text(x, y - 5, notation, color=textcolor,
                ha="left", va="top", fontsize=9, weight="bold")

    plt.show()


def valid_layout_refine(layout_refiner: Role, original_layout: dict, modified_layout: dict, 
                        slide_width: int = 720, slide_height: int = 540,
                        max_aspect_ratio_change: float = 0.2, max_textbox_area_change: float = 0.4,
                        retry: int = 0, max_retries: int = 2) -> dict:
    """
    Validate that the modified layout meets the following criteria:
    1. Image aspect ratios don't change too much
    2. No elements are lost or added
    3. Textbox areas don't change by more than the specified amount
    4. No elements extend outside the slide boundaries
    
    Args:
        layout_refiner: The layout refiner role
        original_layout: The original layout dictionary
        modified_layout: The modified layout dictionary after refinement
        slide_width: Width of the slide in pixels
        slide_height: Height of the slide in pixels
        max_aspect_ratio_change: Maximum allowed aspect ratio change (e.g., 0.2 = 20%)
        max_textbox_area_change: Maximum allowed change in textbox area (e.g., 0.4 = 40%)
        retry: Counter for retry attempts
        max_retries: Maximum number of retries before falling back to original layout
        
    Returns:
        Valid modified layout dictionary or original layout if validation keeps failing
    """
    # Check if we've reached max retries - return original layout if so
    if retry >= max_retries:
        print(f"Reached maximum retries ({max_retries}). Using original layout instead.")
        return original_layout
        
    aspect_ratio_changes = []
    textbox_area_changes = []
    element_changes = []
    boundary_violations = []
    
    # Check if any elements are missing or added
    original_keys = set(original_layout.keys())
    modified_keys = set(modified_layout.keys())

    missing_elements = original_keys - modified_keys
    added_elements = modified_keys - original_keys
    
    if missing_elements:
        element_changes.append(f"Missing elements in modified layout: {', '.join(missing_elements)}")
    
    if added_elements:
        element_changes.append(f"Added elements in modified layout: {', '.join(added_elements)}")
        
    # Check element properties
    for key, value in original_layout.items():
        if key not in modified_layout:
            continue
            
        orig = original_layout[key]
        mod = modified_layout[key]
        
        # Check if element is out of slide bounds
        right_edge = mod['left'] + mod['width']
        bottom_edge = mod['top'] + mod['height']
        
        if mod['left'] < 0 or mod['top'] < 0 or right_edge > slide_width or bottom_edge > slide_height:
            boundary_violations.append({
                'key': key,
                'type': mod.get('type', 'unknown'),
                'position': f"({mod['left']}, {mod['top']}, {right_edge}, {bottom_edge})",
                'slide_bounds': f"(0, 0, {slide_width}, {slide_height})"
            })
        
        # Check image aspect ratios
        if orig.get('type') == 'Picture':
            # Calculate original aspect ratio
            orig_ratio = orig['width'] / orig['height'] if orig['height'] > 0 else 0
            
            # Calculate modified ratio
            mod_ratio = mod['width'] / mod['height'] if mod['height'] > 0 else 0
            
            # Skip if both are zero (would cause division by zero)
            if orig_ratio == 0 or mod_ratio == 0:
                continue
            
            # Calculate percentage change
            ratio_change = abs(orig_ratio - mod_ratio) / orig_ratio
            
            print(f"Image {key} aspect ratio change: {ratio_change:.2f}")
            
            if ratio_change > max_aspect_ratio_change:
                aspect_ratio_changes.append({
                    'key': key,
                    'original_ratio': orig_ratio,
                    'modified_ratio': mod_ratio,
                    'change': ratio_change,
                    'orig_dims': f"{orig['width']}x{orig['height']}",
                    'mod_dims': f"{mod['width']}x{mod['height']}"
                })
        
        # Check textbox areas
        if orig.get('type') == 'TextBox':
            # Calculate original area
            orig_area = orig['width'] * orig['height']
            
            # Calculate modified area
            mod_area = mod['width'] * mod['height']
            
            # Skip if area is zero
            if orig_area == 0:
                continue
            
            # Calculate percentage change
            area_change = abs(orig_area - mod_area) / orig_area
            
            print(f"TextBox {key} area change: {area_change:.2f}")
            
            if area_change > max_textbox_area_change:
                textbox_area_changes.append({
                    'key': key,
                    'original_area': orig_area,
                    'modified_area': mod_area,
                    'change': area_change,
                    'orig_dims': f"{orig['width']}x{orig['height']}",
                    'mod_dims': f"{mod['width']}x{mod['height']}"
                })
    
    print("aspect_ratio_changes:", aspect_ratio_changes)
    print("textbox_area_changes:", textbox_area_changes)
    print("element_changes:", element_changes)
    print("boundary_violations:", boundary_violations)
    
    # If any validation fails, retry
    if aspect_ratio_changes or textbox_area_changes or element_changes or boundary_violations:
        feedback_parts = []
        
        # Add aspect ratio change feedback
        for change in aspect_ratio_changes:
            feedback_parts.append(
                f"Image {change['key']} aspect ratio changed by {change['change']*100:.1f}% "
                f"(from {change['orig_dims']} to {change['mod_dims']}) is not acceptable."
            )
        
        # Add textbox area change feedback
        for change in textbox_area_changes:
            feedback_parts.append(
                f"TextBox {change['key']} area changed by {change['change']*100:.1f}% "
                f"(from original dimensions {change['orig_dims']} to modified dimensions {change['mod_dims']}) is not acceptable."
            )
        
        # Add boundary violation feedback
        for violation in boundary_violations:
            feedback_parts.append(
                f"{violation['type']} {violation['key']} extends outside slide boundaries. "
                f"Position {violation['position']} exceeds slide bounds {violation['slide_bounds']}."
            )
        
        # Add element change feedback
        feedback_parts.extend(element_changes)
        
        feedback = " AND ".join(feedback_parts)
        traceback = f"Layout refinement validation failed: {feedback}"
        print(f"Re-generating layout: {feedback}")
        
        # print("---")
        # # draw_bboxes(original_layout)
        # print("failed layout:")
        # draw_bboxes(modified_layout)
        # print("---")
        
        # # Show retry count in message
        # print(f"Retry attempt {retry+1} of {max_retries}")
        # print("retrying with feedback:", feedback, "and traceback:", traceback)
        
        new_layout = layout_refiner.retry(
            f"Retry to fix the following issues: {feedback}. Ensure all original elements are preserved, "
            f"no new elements are added, image aspect ratios are maintained, and all elements stay within "
            f"the slide boundaries (0,0,{slide_width},{slide_height}).",
            traceback,
            retry + 1
        )
        
        # Validate the new layout
        return valid_layout_refine(layout_refiner, original_layout, new_layout, 
                                    slide_width, slide_height,
                                    max_aspect_ratio_change, max_textbox_area_change,
                                    retry + 1, max_retries)
    
    return modified_layout



def apply_layout_changes(slide: SlidePage, original_layout: dict, target_layout: dict) -> SlidePage:
    """
    Apply layout changes to a slide based on the difference between original and target layouts.
    Assumes that slide elements and original_layout are already matched by indices/keys.
    
    Args:
        slide: The slide to modify
        original_layout: The original layout dictionary
        target_layout: The target layout dictionary with desired changes
        
    Returns:
        Modified slide with applied layout changes
    """
    # Make a complete deep copy of the slide
    modified_slide = deepcopy(slide)
    
    # Process each element in the target layout
    for key, target_elem in target_layout.items():
        if key not in original_layout:
            print(f"Warning: Element {key} not found in original layout. Skipping.")
            continue
            
        # Get the index as an integer
        idx = int(key)
        
        # Skip if index is out of range
        if idx >= len(modified_slide.shapes):
            print(f"Warning: Shape index {idx} out of range (max: {len(modified_slide.shapes)-1}). Skipping.")
            continue
        
        # Get the shape directly by index
        shape = modified_slide.shapes[idx]
        
        # Skip if no change in position or size
        orig_elem = original_layout[key]
        if (orig_elem['left'] == target_elem['left'] and 
            orig_elem['top'] == target_elem['top'] and
            orig_elem['width'] == target_elem['width'] and
            orig_elem['height'] == target_elem['height']):
            continue
            
        # Apply new position and size
        new_left = Pt(target_elem['left'])
        new_top = Pt(target_elem['top'])
        new_width = Pt(target_elem['width'])
        new_height = Pt(target_elem['height'])
        
        # Update the shape bounds in the style dictionary
        shape.style["shape_bounds"]["left"] = new_left
        shape.style["shape_bounds"]["top"] = new_top
        shape.style["shape_bounds"]["width"] = new_width
        shape.style["shape_bounds"]["height"] = new_height
        
        # Different handling based on shape type
        shape_type = original_layout[key].get('type')
        
        # Reset closures to avoid conflicts
        for closure_type in shape._closures:
            shape._closures[closure_type] = []
        
        # For textboxes, use a more explicit approach to ensure the changes take effect
        if shape_type == 'TextBox':
            # Directly set properties using lambda functions
            shape._closures["style"].append(
                Closure(
                    partial(
                        lambda s, l, t, w, h: (
                            setattr(s, 'left', l),
                            setattr(s, 'top', t),
                            setattr(s, 'width', w),
                            setattr(s, 'height', h)
                        )[0],  # Return first element of tuple to make it work
                        new_left, new_top, new_width, new_height
                    ),
                    -1
                )
            )
        else:
            # The simple direct approach for position setting
            shape._closures["style"].append(
                Closure(
                    partial(
                        lambda s, l, t: setattr(s, 'left', l) or setattr(s, 'top', t),
                        new_left, new_top
                    ),
                    -1
                )
            )
            
            # The simple direct approach for size setting
            shape._closures["style"].append(
                Closure(
                    partial(
                        lambda s, w, h: setattr(s, 'width', w) or setattr(s, 'height', h),
                        new_width, new_height
                    ),
                    -1
                )
            )
        
        print(f"Updated shape {idx} ({shape.__class__.__name__}): position ({orig_elem['left']}, {orig_elem['top']}) → ({target_elem['left']}, {target_elem['top']}), size {orig_elem['width']}x{orig_elem['height']} → {target_elem['width']}x{target_elem['height']}")
    
    return modified_slide



def save_modified_presentation(original_prs, config, modified_slide, target_slide_idx, output_path):
    """
    Save the modified presentation with the updated slide.
    
    Args:
        original_prs: Original presentation object
        modified_slide: Modified slide to insert
        target_slide_idx: Index of the slide to replace (1-based)
        output_path: Path to save the modified presentation
    """
    # Deep copy the entire presentation
    new_prs = deepcopy(original_prs)
    
    # Replace the target slide with our modified slide
    if target_slide_idx <= len(new_prs.slides):
        new_prs.slides[target_slide_idx - 1] = modified_slide
        print(f"Replaced slide {target_slide_idx} with modified slide")
    else:
        print(f"Warning: Target slide index {target_slide_idx} is out of range (max: {len(new_prs.slides)})")
        new_prs.slides.append(modified_slide)
    
    # Save the presentation
    new_prs.save(output_path)
    
    print(f"Successfully saved modified presentation to {output_path} with {len(new_prs.slides)} slides")
    
    return new_prs
