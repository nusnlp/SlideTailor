import inspect
import os
import re
import traceback
from copy import deepcopy
from dataclasses import dataclass
from enum import Enum
from functools import partial
from typing import Union

import PIL
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.shapes.base import BaseShape
from pptx.util import Pt

from presentation import Closure, Picture, ShapeElement, SlidePage, TextFrame
from utils import runs_merge


@dataclass
class HistoryMark:
    """
    Mark the execution status of the API call, comment and a line of code.
    """

    API_CALL_ERROR = "api_call_error"
    API_CALL_CORRECT = "api_call_correct"
    COMMENT_CORRECT = "comment_correct"
    COMMENT_ERROR = "comment_error"
    CODE_RUN_ERROR = "code_run_error"
    CODE_RUN_CORRECT = "code_run_correct"



class CodeExecutor:
    """
    Execute code actions and manage API call history, and providing error feedback.
    """

    def __init__(self, retry_times: int):
        """
        Initialize the CodeExecutor.

        Args:
            retry_times (int): The number of times to retry failed actions.
        """
        self.api_history = []
        self.command_history = []
        self.code_history = []
        self.retry_times = retry_times
        self.registered_functions = API_TYPES.all_funcs()
        self.function_regex = re.compile(r"^[a-z]+_[a-z_]+\(.+\)")

    def get_apis_docs(self, funcs: list[callable], show_example: bool = True) -> str:
        """
        Get the documentation for a list of API functions.

        Args:
            funcs (list[callable]): A list of functions to document.
            show_example (bool): Whether to show examples in the documentation.

        Returns:
            str: The formatted API documentation.
        """
        api_doc = []
        for func in funcs:
            sig = inspect.signature(func)
            params = []
            for name, param in sig.parameters.items():
                if name == "slide":
                    continue
                param_str = name
                if param.annotation != inspect.Parameter.empty:
                    param_str += f": {param.annotation.__name__}"
                if param.default != inspect.Parameter.empty:
                    param_str += f" = {repr(param.default)}"
                params.append(param_str)
            signature = f"def {func.__name__}({', '.join(params)})"
            if not show_example:
                api_doc.append(signature)
                continue
            doc = inspect.getdoc(func)
            if doc is not None:
                signature += f"\n\t{doc}"
            api_doc.append(signature)
        return "\n\n".join(api_doc)

    def execute_actions(
        self, actions: str, edit_slide: SlidePage, found_code: bool = False
    ) -> Union[tuple[str, str], None]:
        """
        Execute a series of actions on a slide.

        Args:
            actions (str): The actions to execute.
            edit_slide (SlidePage): The slide to edit.
            found_code (bool): Whether code was found in the actions.

        Returns:
            tuple: The API lines and traceback if an error occurs.
            None: If no error occurs.
        """
        api_calls = actions.strip().split("\n")
        self.api_history.append(
            [HistoryMark.API_CALL_ERROR, edit_slide.slide_idx, actions]
        )
        for line_idx, line in enumerate(api_calls):
            try:
                if line_idx == len(api_calls) - 1 and not found_code:
                    raise ValueError(
                        "No code block found in the output, please output the api calls without any prefix."
                    )
                if line.startswith("def"):
                    raise PermissionError("The function definition were not allowed.")
                if line.startswith("#"):
                    if len(self.command_history) != 0:
                        self.command_history[-1][0] = HistoryMark.COMMENT_CORRECT
                    self.command_history.append([HistoryMark.COMMENT_ERROR, line, None])
                    continue
                if not self.function_regex.match(line):
                    continue
                found_code = True
                func = line.split("(")[0]
                if func not in self.registered_functions:
                    raise NameError(f"The function {func} is not defined.")
                # only one of clone and del can be used in a row
                if func.startswith("clone") or func.startswith("del"):
                    tag = func.split("_")[0]
                    if (
                        self.command_history[-1][-1] == None
                        or self.command_history[-1][-1] == tag
                    ):
                        self.command_history[-1][-1] = tag
                    else:
                        raise ValueError(
                            "Invalid command: Both 'clone_paragraph' and 'del_paragraph'/'del_image' are used within a single command. "
                            "Each command must only perform one of these operations based on the quantity_change."
                        )
                self.code_history.append([HistoryMark.CODE_RUN_ERROR, line, None])
                partial_func = partial(self.registered_functions[func], edit_slide)
                eval(line, {}, {func: partial_func})
                self.code_history[-1][0] = HistoryMark.CODE_RUN_CORRECT
            except:
                trace_msg = traceback.format_exc()
                if len(self.code_history) != 0:
                    self.code_history[-1][-1] = trace_msg
                api_lines = (
                    "\n".join(api_calls[: line_idx - 1])
                    + f"\n--> Error Line: {line}\n"
                    + "\n".join(api_calls[line_idx + 1 :])
                )
                return api_lines, trace_msg
        if len(self.command_history) != 0:
            self.command_history[-1][0] = HistoryMark.COMMENT_CORRECT
        self.api_history[-1][0] = HistoryMark.API_CALL_CORRECT


# supporting functions
def element_index(slide: SlidePage, element_id: int) -> ShapeElement:
    """
    Find the an element in a slide.

    Args:
        slide (SlidePage): The slide
        element_id (int): The ID of the element.

    Returns:
        ShapeElement: The shape corresponding to the element ID.

    Raises:
        IndexError: If the element is not found.
    """
    for shape in slide:
        if shape.shape_idx == element_id:
            return shape
    raise IndexError(f"Cannot find element {element_id}, is it deleted or not exist?")


def replace_para(paragraph_id: int, new_text: str, shape: BaseShape):
    """
    Replace the text of a paragraph in a shape.
    """
    para = shape.text_frame.paragraphs[paragraph_id]
    runs_merge(para).text = new_text


def clone_para(paragraph_id: int, shape: BaseShape):
    """
    Clone a paragraph in a shape.
    """
    para = shape.text_frame.paragraphs[paragraph_id]
    shape.text_frame.paragraphs[-1]._element.addnext(parse_xml(para._element.xml))


def clone_img(shape: BaseShape, img_path: str, **style_bounds):
    """
    Clone an image in a slide by adding a new picture with the same properties.
    """
    slide = shape.part.slide
    new_shape = slide.shapes.add_picture(img_path, **style_bounds)
    new_shape.name = shape.name
    if hasattr(shape, 'image'):
        if hasattr(new_shape, 'image'):
            for attr in ['crop_bottom', 'crop_top', 'crop_left', 'crop_right']:
                if hasattr(shape.image, attr):
                    setattr(new_shape.image, attr, getattr(shape.image, attr))
    
    if hasattr(shape, 'line') and hasattr(new_shape, 'line'):
        for attr in ['width', 'dash_style']:
            if hasattr(shape.line, attr):
                setattr(new_shape.line, attr, getattr(shape.line, attr))
    
    if hasattr(shape, 'rotation') and hasattr(new_shape, 'rotation'):
        new_shape.rotation = shape.rotation
    
    return new_shape


def del_para(paragraph_id: int, shape: BaseShape):
    """
    Delete a paragraph from a shape.
    """
    para = shape.text_frame.paragraphs[paragraph_id]
    para._element.getparent().remove(para._element)


# api functions
def del_paragraph(slide: SlidePage, div_id: int, paragraph_id: int):
    """
    Delete a paragraph from a slide.

    Args:
        slide (SlidePage): The slide containing the paragraph.
        div_id (int): The ID of the division containing the paragraph.
        paragraph_id (int): The ID of the paragraph to delete.

    Raises:
        IndexError: If the paragraph is not found.
    """
    shape = element_index(slide, div_id)
    assert (
        shape.text_frame.is_textframe
    ), "The element does not have a text frame, please check the element id and type of element."
    for para in shape.text_frame.paragraphs:
        if para.idx == paragraph_id:
            shape.text_frame.paragraphs.remove(para)
            shape._closures["delete"].append(
                Closure(partial(del_para, para.real_idx), para.real_idx)
            )
            return
    else:
        raise IndexError(
            f"Cannot find the paragraph {paragraph_id} of the element {div_id},"
            "may refer to a non-existed paragraph or you haven't cloned enough paragraphs beforehand."
        )


def del_image(slide: SlidePage, figure_id: int):
    """
    Delete an image from a slide.

    Args:
        slide (SlidePage): The slide containing the image.
        figure_id (int): The ID of the image to delete.
    """
    shape = element_index(slide, figure_id)
    # import pdb; pdb.set_trace()
    # assert isinstance(shape, Picture), "The element is not a Picture."
    try:
        slide.shapes.remove(shape)
    except:
        pass


def replace_paragraph(slide: SlidePage, div_id: int, paragraph_id: int, text: str):
    """
    Replace the text of a paragraph in a slide.

    Args:
        slide (SlidePage): The slide containing the paragraph.
        div_id (int): The ID of the division containing the paragraph.
        paragraph_id (int): The ID of the paragraph to replace.
        text (str): The new text to replace with.

    Raises:
        IndexError: If the paragraph is not found.
    """
    shape = element_index(slide, div_id)
    assert (
        shape.text_frame.is_textframe
    ), "The element does not have a text frame, please check the element id and type of element."
    for para in shape.text_frame.paragraphs:
        if para.idx == paragraph_id:
            para.text = text
            shape._closures["replace"].append(
                Closure(
                    partial(replace_para, para.real_idx, text),
                    para.real_idx,
                )
            )
            return
    else:
        raise IndexError(
            f"Cannot find the paragraph {paragraph_id} of the element {div_id},"
            "Please: "
            "1. check if you refer to a non-existed paragraph."
            "2. check if you already deleted it."
        )


def replace_image(slide: SlidePage, img_id: int, image_path: str):
    """
    Replace an image in a slide.

    Args:
        slide (SlidePage): The slide containing the image.
        img_id (int): The ID of the image to replace.
        image_path (str): The path to the new image.

    Raises:
        ValueError: If the image path does not exist.
    """
    if not os.path.exists(image_path):
        raise ValueError(
            f"The image {image_path} does not exist, consider use del_image if image_path in the given command is faked"
        )
    shape = element_index(slide, img_id)
    assert isinstance(shape, Picture), "The element is not a Picture."
    img_size = PIL.Image.open(image_path).size
    r = min(shape.width / img_size[0], shape.height / img_size[1])
    new_width = int(img_size[0] * r)
    new_height = int(img_size[1] * r)
    shape.width = Pt(new_width)
    shape.height = Pt(new_height)
    shape.img_path = image_path
    
    # import pdb; pdb.set_trace()




def clone_image(slide: SlidePage, img_id: int):
    """
    Clone an image in a slide.

    Args:
        slide (SlidePage): The slide containing the image.
        img_id (int): The ID of the image to clone.

    Raises:
        IndexError: If the image is not found.
        ValueError: If the element is not an image.

    Note: The cloned image will have an image_id one greater than the current maximum in the slide.
    """
    shape = element_index(slide, img_id)
    assert isinstance(shape, Picture), "The element is not a Picture."
    
    # Find max shape index in the slide
    max_idx = max([s.shape_idx for s in slide])
    
    # Create a deep copy of the image shape
    new_shape = deepcopy(shape)
    new_shape.shape_idx = max_idx + 1
    
    new_shape.style["shape_bounds"]["width"] = shape.style["shape_bounds"]["width"]
    new_shape.style["shape_bounds"]["height"] = shape.style["shape_bounds"]["height"]
    new_shape.style["shape_bounds"]["left"] = int(shape.style["shape_bounds"]["left"] * 1.2)
    new_shape.style["shape_bounds"]["top"] = int(shape.style["shape_bounds"]["top"] * 1.2)
    
    # Add the cloned shape to the slide shapes
    slide.shapes.append(new_shape)
    
    # Add a closure to handle the actual image cloning in PowerPoint
    # shape._closures["clone"].append(
    #     Closure(
    #         partial(
    #             clone_img, 
    #             img_path=shape.img_path, 
    #             **shape.style["shape_bounds"]
    #         ),
    #         -1  # Not paragraph specific
    #     )
    # )
    
    # import pdb; pdb.set_trace()
    


def add_image(slide: SlidePage, image_path: str,
            left: float | None = 0, top: float | None = 0,
            width: float | None = None, height: float | None = None):
    """
    Add a new image to a slide at the specified position and dimensions.

    Args:
        slide (SlidePage): The slide to add the image to.
        image_path (str): Path to the image file.
        left (float, optional): Left position in points. Defaults to 0.
        top (float, optional): Top position in points. Defaults to 0.
        width (float, optional): Width in points. Defaults to original image width if None.
        height (float, optional): Height in points. Defaults to original image height if None.

    Raises:
        ValueError: If the image path does not exist.
        ValueError: If the dimensions are invalid.
    """
    if not os.path.exists(image_path):
        raise ValueError(f"The image {image_path} does not exist")
    
    # Define slide boundaries with safety margin
    MARGIN = 20  # 20pt margin to ensure the image stays fully within the slide
    max_width = slide.slide_width - (2 * MARGIN)
    max_height = slide.slide_height - (2 * MARGIN)
    
    # Get original image dimensions
    img_size = PIL.Image.open(image_path).size
    original_width, original_height = img_size[0], img_size[1]
    original_aspect_ratio = original_width / original_height
    
    # Calculate dimensions if not provided
    if width is None and height is None:
        # Use original dimensions
        width = original_width
        height = original_height
    elif width is None:
        # Calculate width based on height to maintain aspect ratio
        width = height * original_aspect_ratio
    elif height is None:
        # Calculate height based on width to maintain aspect ratio
        height = width / original_aspect_ratio
    
    # Ensure image fits within slide boundaries
    if width > max_width or height > max_height:
        # Calculate scale factors for width and height
        width_scale = max_width / width
        height_scale = max_height / height
        
        # Use the smallest scale factor to ensure image fits within slide
        scale = min(width_scale, height_scale)
        
        # Apply scaling while maintaining aspect ratio
        width = width * scale
        height = height * scale
    
    # Convert to Pt for consistency
    width_pt = Pt(width)
    height_pt = Pt(height)
    
    # Process position values
    if left is None:
        # Center horizontally
        left = (slide.slide_width - width) / 2
    else:
        # Ensure left position keeps image inside slide
        left = max(MARGIN, min(slide.slide_width - width - MARGIN, left))
        
    if top is None:
        # Center vertically
        top = (slide.slide_height - height) / 2
    else:
        # Ensure top position keeps image inside slide
        top = max(MARGIN, min(slide.slide_height - height - MARGIN, top))
    
    # Convert to Pt
    left_pt = Pt(left)
    top_pt = Pt(top)
    
    # Find max shape index in the slide
    max_idx = max([s.shape_idx for s in slide]) if slide.shapes else 0
    
    # Create style dictionary
    style = {
        "shape_bounds": {
            "width": width_pt,
            "height": height_pt,
            "left": left_pt,
            "top": top_pt,
        },
        "shape_type": "picture",
        "rotation": 0,
        "fill": None,
        "line": None,
        "img_style": {
            "crop_bottom": 0,
            "crop_top": 0, 
            "crop_left": 0,
            "crop_right": 0
        }
    }
    
    # Create a text frame (empty for images)
    base_shape = BaseShape(shape_elm=None, parent=slide)
    text_frame = TextFrame(base_shape, 0)
    
    # Create the new Picture shape
    new_shape = Picture(
        slide.slide_idx,
        max_idx + 1,
        style,
        [image_path, os.path.basename(image_path), ""],  # [img_path, name, caption]
        text_frame,
        slide.slide_width * slide.slide_height,
        level=0
    )
    
    # Add to slide's shapes
    slide.shapes.append(new_shape)


def detect_overlap(bounds1, bounds2, buffer=0):
    """
    Detect if two bounding boxes overlap.
    
    Args:
        bounds1 (dict): First bounding box with 'left', 'top', 'width', 'height' keys
        bounds2 (dict): Second bounding box with 'left', 'top', 'width', 'height' keys
        buffer (float): Optional buffer space to add around elements (in points)
        
    Returns:
        bool: True if the bounding boxes overlap, False otherwise
    """
    # Extract coordinates with buffer
    left1 = bounds1['left'].pt if hasattr(bounds1['left'], 'pt') else bounds1['left']
    top1 = bounds1['top'].pt if hasattr(bounds1['top'], 'pt') else bounds1['top']
    width1 = bounds1['width'].pt if hasattr(bounds1['width'], 'pt') else bounds1['width']
    height1 = bounds1['height'].pt if hasattr(bounds1['height'], 'pt') else bounds1['height']
    
    left2 = bounds2['left'].pt if hasattr(bounds2['left'], 'pt') else bounds2['left']
    top2 = bounds2['top'].pt if hasattr(bounds2['top'], 'pt') else bounds2['top']
    width2 = bounds2['width'].pt if hasattr(bounds2['width'], 'pt') else bounds2['width']
    height2 = bounds2['height'].pt if hasattr(bounds2['height'], 'pt') else bounds2['height']
    
    # Add buffer
    left1 -= buffer
    top1 -= buffer
    width1 += 2 * buffer
    height1 += 2 * buffer
    
    # Check for overlap
    if (left1 + width1 <= left2 or  # 1 is left of 2
        left2 + width2 <= left1 or  # 2 is left of 1
        top1 + height1 <= top2 or   # 1 is above 2
        top2 + height2 <= top1):    # 2 is above 1
        return False
    return True


def get_overlap_area(bounds1, bounds2):
    """
    Calculate the area of overlap between two bounding boxes.
    
    Args:
        bounds1 (dict): First bounding box with 'left', 'top', 'width', 'height' keys
        bounds2 (dict): Second bounding box with 'left', 'top', 'width', 'height' keys
        
    Returns:
        float: Area of overlap in square points, 0 if no overlap
    """
    if not detect_overlap(bounds1, bounds2):
        return 0
    
    # Extract coordinates
    left1 = bounds1['left'].pt if hasattr(bounds1['left'], 'pt') else bounds1['left']
    top1 = bounds1['top'].pt if hasattr(bounds1['top'], 'pt') else bounds1['top']
    right1 = left1 + (bounds1['width'].pt if hasattr(bounds1['width'], 'pt') else bounds1['width'])
    bottom1 = top1 + (bounds1['height'].pt if hasattr(bounds1['height'], 'pt') else bounds1['height'])
    
    left2 = bounds2['left'].pt if hasattr(bounds2['left'], 'pt') else bounds2['left']
    top2 = bounds2['top'].pt if hasattr(bounds2['top'], 'pt') else bounds2['top']
    right2 = left2 + (bounds2['width'].pt if hasattr(bounds2['width'], 'pt') else bounds2['width'])
    bottom2 = top2 + (bounds2['height'].pt if hasattr(bounds2['height'], 'pt') else bounds2['height'])
    
    # Calculate overlap dimensions
    overlap_width = max(0, min(right1, right2) - max(left1, left2))
    overlap_height = max(0, min(bottom1, bottom2) - max(top1, top2))
    
    return overlap_width * overlap_height


def find_optimal_position(shape, other_shapes, slide, max_attempts=20):
    """
    Find an optimal position for a shape to avoid overlaps with other shapes.
    
    Args:
        shape (ShapeElement): The shape to position
        other_shapes (list): List of other shapes to avoid
        slide (SlidePage): The slide containing the shapes
        max_attempts (int): Maximum number of attempts to find a non-overlapping position
        
    Returns:
        dict: New bounds with 'left' and 'top' keys (in points)
    """
    MARGIN = 20  # Margin from slide edges
    bounds = shape.style["shape_bounds"]
    
    # Extract values and ensure they're in pts
    width = bounds['width'].pt if hasattr(bounds['width'], 'pt') else bounds['width']
    height = bounds['height'].pt if hasattr(bounds['height'], 'pt') else bounds['height']
    left = bounds['left'].pt if hasattr(bounds['left'], 'pt') else bounds['left']
    top = bounds['top'].pt if hasattr(bounds['top'], 'pt') else bounds['top']
    
    # Try grid positions
    grid_size = min(slide.slide_width, slide.slide_height) / 10
    
    best_position = {'left': left, 'top': top}
    lowest_overlap = float('inf')
    
    for attempt in range(max_attempts):
        # Try different positions in a spiral pattern or grid
        if attempt == 0:
            # Try center position first
            test_left = (slide.slide_width - width) / 2
            test_top = (slide.slide_height - height) / 2
        else:
            # Grid positions with some randomization
            grid_x = attempt % 5
            grid_y = attempt // 5
            test_left = MARGIN + grid_x * (slide.slide_width - width - 2*MARGIN) / 4
            test_top = MARGIN + grid_y * (slide.slide_height - height - 2*MARGIN) / 4
            
            # Add some randomization to avoid exact grid alignment
            import random
            test_left += random.uniform(-grid_size/4, grid_size/4)
            test_top += random.uniform(-grid_size/4, grid_size/4)
        
        # Ensure within slide boundaries
        test_left = max(MARGIN, min(slide.slide_width - width - MARGIN, test_left))
        test_top = max(MARGIN, min(slide.slide_height - height - MARGIN, test_top))
        
        # Create test bounds
        test_bounds = {
            'left': test_left,
            'top': test_top,
            'width': width,
            'height': height
        }
        
        # Check overlap with other shapes
        total_overlap = 0
        for other in other_shapes:
            other_bounds = other.style["shape_bounds"]
            overlap = get_overlap_area(test_bounds, other_bounds)
            total_overlap += overlap
        
        # Update best position if better
        if total_overlap < lowest_overlap:
            lowest_overlap = total_overlap
            best_position = {'left': test_left, 'top': test_top}
            
            # If no overlap, we can stop searching
            if total_overlap == 0:
                break
    
    return best_position


def auto_rearrange_elements(slide: SlidePage, buffer: float = 10, allow_resize: bool = False, resize_factor: float = 0.9):
    """
    Automatically rearrange image elements on a slide to reduce overlap.
    Ignores non-image elements like TextBox and FreeShape.
    
    Args:
        slide (SlidePage): The slide containing elements to rearrange
        buffer (float): Buffer space to add around elements (in points)
        allow_resize (bool): Whether to allow resizing images if needed to avoid overlaps
        resize_factor (float): Factor to resize images by (between 0 and 1) if resizing is allowed
        
    Returns:
        int: Number of elements repositioned
    """
    if not slide.shapes or len(slide.shapes) <= 1:
        return 0  # Nothing to rearrange
    
    # Filter to only include Picture objects
    image_shapes = [shape for shape in slide.shapes if isinstance(shape, Picture)]
    if len(image_shapes) <= 1:
        return 0  # Nothing to rearrange
    
    # Get non-image elements (to avoid overlapping with them)
    non_image_shapes = [shape for shape in slide.shapes if not isinstance(shape, Picture)]
    
    # Define slide boundaries with safety margin
    MARGIN = 20  # Margin from slide edges
    
    # Sort images by area (larger first)
    images_by_area = sorted(
        image_shapes, 
        key=lambda s: (
            s.width * s.height if hasattr(s, 'width') and hasattr(s, 'height') 
            else s.style["shape_bounds"]["width"].pt * s.style["shape_bounds"]["height"].pt
        ),
        reverse=True
    )
    
    # Track images that have been positioned
    positioned_images = []
    repositioned_count = 0
    
    # Grid-based placement to control positioning
    grid_columns = 2  # Number of columns in the placement grid
    grid_rows = (len(images_by_area) + grid_columns - 1) // grid_columns  # Ceiling division
    
    # Available grid positions
    grid_positions = []
    
    # Calculate width and height available for each grid cell
    grid_width = (slide.slide_width - 2*MARGIN) / grid_columns
    grid_height = (slide.slide_height - 2*MARGIN) / grid_rows
    
    # Generate grid positions
    for row in range(grid_rows):
        for col in range(grid_columns):
            grid_positions.append({
                'left': MARGIN + col * grid_width,
                'top': MARGIN + row * grid_height,
                'width': grid_width - buffer,
                'height': grid_height - buffer
            })
    
    # Start with the largest image (in its original position)
    positioned_images.append(images_by_area[0])
    
    # Position remaining images
    for i, image in enumerate(images_by_area[1:], 1):
        current_bounds = image.style["shape_bounds"]
        
        # Extract current width and height (in points)
        width = current_bounds['width'].pt if hasattr(current_bounds['width'], 'pt') else current_bounds['width']
        height = current_bounds['height'].pt if hasattr(current_bounds['height'], 'pt') else current_bounds['height']
        
        # Check if current position causes overlap
        overlaps = False
        for positioned in positioned_images + non_image_shapes:
            if detect_overlap(current_bounds, positioned.style["shape_bounds"], buffer):
                overlaps = True
                break
        
        if overlaps or i >= len(grid_positions):
            # Get the current grid position or use the best available one
            if i < len(grid_positions):
                grid_pos = grid_positions[i]
            else:
                # Find the first available grid position
                for pos in grid_positions:
                    # Use this position
                    grid_pos = pos
                    break
                else:
                    # If all grid positions are used, use the last one
                    grid_pos = grid_positions[-1]
            
            # If resizing is allowed and the image is too large for the grid cell
            if allow_resize and (width > grid_pos['width'] or height > grid_pos['height']):
                # Calculate scale factors
                width_scale = grid_pos['width'] / width
                height_scale = grid_pos['height'] / height
                
                # Use the smaller scale factor to maintain aspect ratio
                scale = min(width_scale, height_scale, resize_factor)
                
                # Apply scaling
                new_width = width * scale
                new_height = height * scale
                
                # Convert to Pt
                new_width_pt = Pt(new_width)
                new_height_pt = Pt(new_height)
                
                # Update width and height
                image.style["shape_bounds"]["width"] = new_width_pt
                image.style["shape_bounds"]["height"] = new_height_pt
                
                # Add a closure to update the image size in PowerPoint
                image._closures["style"].append(
                    Closure(
                        partial(
                            lambda s, width, height: setattr(s, 'width', width) or setattr(s, 'height', height),
                            new_width_pt, new_height_pt
                        ),
                        -1  # Not paragraph specific
                    )
                )
                
                # Update width and height for position calculation
                width = new_width
                height = new_height
            
            # Calculate new position (center in grid cell)
            new_left = grid_pos['left'] + (grid_pos['width'] - width) / 2
            new_top = grid_pos['top'] + (grid_pos['height'] - height) / 2
            
            # Ensure within slide boundaries
            new_left = max(MARGIN, min(slide.slide_width - width - MARGIN, new_left))
            new_top = max(MARGIN, min(slide.slide_height - height - MARGIN, new_top))
            
            # Convert to Pt
            new_left_pt = Pt(new_left)
            new_top_pt = Pt(new_top)
            
            # Apply new position
            image.style["shape_bounds"]["left"] = new_left_pt
            image.style["shape_bounds"]["top"] = new_top_pt
            
            # Add a closure to update the image position in PowerPoint
            image._closures["style"].append(
                Closure(
                    partial(
                        lambda s, left, top: setattr(s, 'left', left) or setattr(s, 'top', top),
                        new_left_pt, new_top_pt
                    ),
                    -1  # Not paragraph specific
                )
            )
            
            repositioned_count += 1
        
        # Add to positioned images list
        positioned_images.append(image)
    
    return repositioned_count


def clone_paragraph(slide: SlidePage, div_id: int, paragraph_id: int):
    """
    Clone a paragraph in a slide.

    Args:
        slide (SlidePage): The slide containing the paragraph.
        div_id (int): The ID of the division containing the paragraph.
        paragraph_id (int): The ID of the paragraph to clone.

    Raises:
        IndexError: If the paragraph is not found.

    Mention: the cloned paragraph will have a paragraph_id one greater than the current maximum in the parent element.
    """
    shape = element_index(slide, div_id)
    assert (
        shape.text_frame.is_textframe
    ), "The element does not have a text frame, please check the element id and type of element."
    max_idx = max([para.idx for para in shape.text_frame.paragraphs])
    for para in shape.text_frame.paragraphs:
        if para.idx != paragraph_id:
            continue
        shape.text_frame.paragraphs.append(deepcopy(para))
        shape.text_frame.paragraphs[-1].idx = max_idx + 1
        shape.text_frame.paragraphs[-1].real_idx = len(shape.text_frame.paragraphs) - 1
        shape._closures["clone"].append(
            Closure(
                partial(clone_para, para.real_idx),
                para.real_idx,
            )
        )
        return
    raise IndexError(
        f"Cannot find the paragraph {paragraph_id} of the element {div_id}, may refer to a non-existed paragraph."
    )


class API_TYPES(Enum):
    Agent = [
        replace_image,
        del_image,
        clone_paragraph,
        replace_paragraph,
        del_paragraph,
        # clone_image,
        # add_image,
        # auto_rearrange_elements,
    ]

    @classmethod
    def all_funcs(cls) -> dict[str, callable]:
        funcs = {}
        for attr in dir(cls):
            if attr.startswith("__"):
                continue
            funcs |= {func.__name__: func for func in getattr(cls, attr).value}
        return funcs


if __name__ == "__main__":
    print(CodeExecutor(0).get_apis_docs(API_TYPES.Agent.value))
